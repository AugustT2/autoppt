"""
生成两个对比模板（均为示例说明页 + 一页「产品页」骨架）：
1) template_placeholder_comparison.pptx — 使用内置「比较」版式的幻灯片占位符
2) template_named_shapes_blank.pptx — 使用空白版式 + 命名形状（便于程序按 name 查找）

运行: python build_template_demos.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

DIR = r"d:\IdeaProjects\autoppt"
OUT_PLACEHOLDER = rf"{DIR}\template_placeholder_comparison.pptx"
OUT_NAMED = rf"{DIR}\template_named_shapes_blank.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_cnvpr_name(shape, new_name: str) -> None:
    el = shape._element
    for tag in ("nvSpPr", "nvPicPr", "nvCxnSpPr", "nvGraphicFramePr", "nvGrpSpPr"):
        block = getattr(el, tag, None)
        if block is not None and getattr(block, "cNvPr", None) is not None:
            block.cNvPr.set("name", new_name)
            return
    raise ValueError(f"无法设置形状名称: {type(shape)}")


def _fill_placeholder_text(shape, lines: list[tuple[str, bool]], *, size=Pt(11)) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.name = "微软雅黑"
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


def build_placeholder_template() -> None:
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    # 说明页
    tslide = prs.slides.add_slide(prs.slide_layouts[0])
    tslide.shapes.title.text = "模板 A：幻灯片占位符（Comparison 版式）"
    sub = tslide.placeholders[1]
    _fill_placeholder_text(
        sub,
        [
            ("本文件第 2 页使用 PowerPoint 内置「比较」版式，灰色框是版式自带的占位符。", False),
            ("", False),
            ("自动化写法（python-pptx 示例思路）：slide.placeholders[0] 写标题；[1]～[4] 按 idx 填文字/插入表格/图表。", False),
            ("", False),
            ("特点：对象与母版版式绑定；新建同版式幻灯片时占位符位置一致。", False),
            ("注意：占位符类型/编号随版式变化，换版式要改代码下标。", True),
        ],
        size=Pt(14),
    )

    # 产品骨架：Comparison = Title + 左文 + 左对象 + 右文 + 右对象
    slide = prs.slides.add_slide(prs.slide_layouts[4])

    _fill_placeholder_text(
        slide.placeholders[0],
        [
            ("【占位符 idx=0 · TITLE】产品页主标题", True),
            ("示例：「基金业绩说明｜蓝海稳健增长混合 A」", False),
        ],
        size=Pt(20),
    )
    ph1 = slide.placeholders[1]
    _fill_placeholder_text(
        ph1,
        [
            ("【idx=1 · TEXT】左侧小标题区", True),
            ("可写：基金基本情况、费率等章节标题。", False),
        ],
    )
    ph2 = slide.placeholders[2]
    _fill_placeholder_text(
        ph2,
        [
            ("【idx=2 · OBJECT】左侧主体区", True),
            ("通常放多行简介文字，或由程序在此占位符内插入表格。", False),
            ("（表格/图表也可通过占位符区域插入，取决于实现方式）", False),
        ],
    )
    ph3 = slide.placeholders[3]
    _fill_placeholder_text(
        ph3,
        [
            ("【idx=3 · TEXT】右侧小标题区", True),
            ("可写：业绩排名、走势图说明等。", False),
        ],
    )
    ph4 = slide.placeholders[4]
    _fill_placeholder_text(
        ph4,
        [
            ("【idx=4 · OBJECT】右侧主体区", True),
            ("通常放第二张图或补充表格。", False),
        ],
    )

    # 浅色提示底（可选，画在占位符之上会破坏占位符——这里用边框色区分：略）
    for ph in (ph1, ph2, ph3, ph4):
        try:
            ph.fill.solid()
            ph.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFF)
        except Exception:
            pass

    prs.save(OUT_PLACEHOLDER)
    print("Saved:", OUT_PLACEHOLDER)


def build_named_shapes_template() -> None:
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    tslide = prs.slides.add_slide(prs.slide_layouts[0])
    tslide.shapes.title.text = "模板 B：命名形状（Blank 空白版式）"
    sub = tslide.placeholders[1]
    _fill_placeholder_text(
        sub,
        [
            ("本文件第 2 页为空白版式，每个可编辑区域都是普通形状，并在选择窗格中命名为英文常量。", False),
            ("", False),
            ("自动化写法：遍历 slide.shapes，按 shape.name == \"TBL_RANKING\" 等查找后填数/换图。", False),
            ("", False),
            ("特点：版式自由，接近「照着公司范例画一页」；名字由模板约定，与占位符 idx 无关。", False),
            ("请在 PowerPoint「开始 → 选择 → 选择窗格」中核对名字。", True),
        ],
        size=Pt(14),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    def add_labeled_box(name: str, left, top, width, height, title: str, body: str) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        _set_cnvpr_name(box, name)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.bold = True
        p0.font.size = Pt(11)
        p0.font.name = "微软雅黑"
        p0.font.color.rgb = RGBColor(0x2E, 0x40, 0x53)
        p1 = tf.add_paragraph()
        p1.text = body
        p1.font.size = Pt(10)
        p1.font.name = "微软雅黑"
        p1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p1.space_before = Pt(4)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFF)
        box.line.color.rgb = RGBColor(0x99, 0xB4, 0xCC)
        box.line.width = Pt(0.75)

    y0 = Inches(0.72)
    mid = Inches(3.72)
    lw = Inches(6.15)
    rw = Inches(6.35)
    rx = Inches(6.75)

    add_labeled_box(
        "TITLE_PAGE",
        Inches(0.4),
        Inches(0.12),
        SLIDE_W - Inches(0.8),
        Inches(0.52),
        "TITLE_PAGE（大标题）",
        "程序替换为：基金业绩说明 + 产品名",
    )

    add_labeled_box(
        "LBL_FUND_SECTION",
        Inches(0.45),
        y0,
        lw,
        Inches(0.28),
        "LBL_FUND_SECTION",
        "小节标题占位，可改为静态「基金基本情况」",
    )
    add_labeled_box(
        "TXT_FUND_PROFILE",
        Inches(0.45),
        y0 + Inches(0.32),
        lw,
        Inches(2.45),
        "TXT_FUND_PROFILE（多行文字）",
        "程序写入：类型、成立日、规模、经理、基准、风险等级等。",
    )

    add_labeled_box(
        "LBL_RANK_SECTION",
        Inches(0.45),
        mid,
        lw,
        Inches(0.28),
        "LBL_RANK_SECTION",
        "小节标题，如「收益率排名」",
    )

    rows, cols = 4, 4
    tbl_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.45),
        mid + Inches(0.32),
        lw - Inches(0.1),
        Inches(1.38),
    )
    _set_cnvpr_name(tbl_shape, "TBL_RANKING")
    table = tbl_shape.table
    hdr = ["区间", "收益率", "同类排名", "分位数"]
    for c, h in enumerate(hdr):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(9)
            p.font.name = "微软雅黑"
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = "（数据）"

    add_labeled_box(
        "LBL_ALLOC_SECTION",
        rx,
        y0,
        rw,
        Inches(0.28),
        "LBL_ALLOC_SECTION",
        "如：过去4季度资产配置（%）",
    )
    cd1 = CategoryChartData()
    cd1.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd1.add_series("股票", (60, 62, 58, 65))
    cd1.add_series("债券", (30, 28, 32, 27))
    ch1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        rx,
        y0 + Inches(0.32),
        rw,
        Inches(2.45),
        cd1,
    )
    _set_cnvpr_name(ch1, "CHART_ALLOCATION")
    ch1.chart.has_legend = True
    ch1.chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    add_labeled_box(
        "LBL_NAV_SECTION",
        rx,
        mid,
        rw,
        Inches(0.28),
        "LBL_NAV_SECTION",
        "如：累计收益率走势",
    )
    cd2 = CategoryChartData()
    cd2.categories = ["T1", "T2", "T3", "T4", "T5"]
    cd2.add_series("本基金", (0, 2, 4, 5, 8))
    cd2.add_series("基准", (0, 1, 2, 3, 5))
    ch2 = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        rx,
        mid + Inches(0.32),
        rw,
        Inches(2.95),
        cd2,
    )
    _set_cnvpr_name(ch2, "CHART_NAV_SERIES")
    ch2.chart.has_legend = True
    ch2.chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # 分隔线（无业务名，仅装饰）
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.48), y0 - Inches(0.05), Inches(0.02), Inches(6.35))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln.line.fill.background()
    ln2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), mid - Inches(0.22), Inches(12.6), Inches(0.02)
    )
    ln2.fill.solid()
    ln2.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln2.line.fill.background()

    # 角标说明
    cap = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.38))
    cap.text_frame.paragraphs[0].text = (
        "命名形状：TITLE_PAGE · TXT_FUND_PROFILE · TBL_RANKING · CHART_ALLOCATION · CHART_NAV_SERIES "
        "（及 LBL_* 小节标题框，可按需精简）"
    )
    cap.text_frame.paragraphs[0].font.size = Pt(9)
    cap.text_frame.paragraphs[0].font.name = "微软雅黑"
    cap.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUT_NAMED)
    print("Saved:", OUT_NAMED)


def main() -> None:
    build_placeholder_template()
    build_named_shapes_template()


if __name__ == "__main__":
    main()

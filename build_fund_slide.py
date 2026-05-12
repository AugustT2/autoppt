"""Generate a single-slide fund performance PPT (sample data)."""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = r"d:\IdeaProjects\autoppt\fund_performance_overview.pptx"

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title(slide, prs):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), SLIDE_W - Inches(0.8), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "基金业绩说明（示例数据）"
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1C, 0x28, 0x33)
    p.alignment = PP_ALIGN.CENTER


def add_quadrant_label(slide, left, top, text):
    box = slide.shapes.add_textbox(left, top, Inches(4.5), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2E, 0x40, 0x53)


def add_body_text(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(0x29, 0x29, 0x29)
        p.space_after = Pt(4)


def add_ranking_table(slide, left, top, width, height):
    rows, cols = 4, 4
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["区间", "收益率", "同类排名", "分位数"]
    data = [
        ["近一年", "+12.8%", "128 / 856", "前 15%"],
        ["近两年", "+18.6%", "95 / 812", "前 12%"],
        ["近三年", "+22.4%", "76 / 698", "前 11%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(9.5)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2E, 0x40, 0x53)
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5)
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0x29, 0x29, 0x29)
    for col in table.columns:
        col.width = int(width / cols)


def main():
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_title(slide, prs)

    # Quadrant positions (below title)
    y0 = Inches(0.75)
    mid_y = Inches(3.85)
    left_w = Inches(6.35)
    right_x = Inches(6.55)
    right_w = Inches(6.55)
    half_h = Inches(3.0)

    add_quadrant_label(slide, Inches(0.45), y0, "基金基本情况")
    fund_lines = [
        ("基金名称：蓝海稳健增长混合A（示例）", True),
        ("基金类型：偏股混合型基金", False),
        ("成立日期：2019-06-12　　最新规模：58.6 亿元（示例）", False),
        ("基金经理：张明、李悦（示例）", False),
        ("业绩比较基准：沪深300指数收益率×70% + 中债综合指数×30%", False),
        ("风险等级：R3（中风险）　　托管人：示例商业银行", False),
    ]
    add_body_text(slide, Inches(0.45), Inches(1.05), left_w - Inches(0.5), Inches(2.55), fund_lines)

    add_quadrant_label(slide, Inches(0.45), mid_y - Inches(0.15), "收益率排名（同类偏股混合）")
    add_ranking_table(slide, Inches(0.45), mid_y + Inches(0.12), left_w - Inches(0.5), Inches(1.35))

    add_quadrant_label(slide, right_x, y0, "过去4个季度资产类别配置（%）")
    chart_data_alloc = CategoryChartData()
    chart_data_alloc.categories = ["2024Q2", "2024Q3", "2024Q4", "2025Q1"]
    chart_data_alloc.add_series("股票", (68.2, 71.5, 65.8, 72.4))
    chart_data_alloc.add_series("债券", (22.5, 19.8, 24.1, 18.6))
    chart_data_alloc.add_series("现金及其他", (9.3, 8.7, 10.1, 9.0))
    chart_alloc = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        right_x,
        Inches(1.05),
        right_w,
        Inches(2.65),
        chart_data_alloc,
    ).chart
    chart_alloc.has_legend = True
    chart_alloc.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart_alloc.legend.include_in_layout = False

    add_quadrant_label(slide, right_x, mid_y - Inches(0.15), "累计收益率走势（示例，%）")
    chart_data_ret = CategoryChartData()
    chart_data_ret.categories = [
        "2024-05",
        "2024-07",
        "2024-09",
        "2024-11",
        "2025-01",
        "2025-03",
        "2025-05",
    ]
    chart_data_ret.add_series("本基金", (0.0, 3.2, 5.1, 7.8, 6.2, 9.5, 12.8))
    chart_data_ret.add_series("业绩基准", (0.0, 2.1, 3.8, 5.5, 4.2, 6.1, 8.4))
    chart_ret = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        right_x,
        mid_y + Inches(0.12),
        right_w,
        Inches(2.95),
        chart_data_ret,
    ).chart
    chart_ret.has_legend = True
    chart_ret.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart_ret.legend.include_in_layout = False

    # Subtle divider lines
    from pptx.enum.shapes import MSO_SHAPE

    line1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(0.72), Inches(0.02), Inches(6.5)
    )
    line1.fill.solid()
    line1.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line1.line.fill.background()
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), mid_y - Inches(0.22), Inches(12.6), Inches(0.02)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line2.line.fill.background()

    prs.save(OUT)
    print("Saved:", OUT)


if __name__ == "__main__":
    main()
